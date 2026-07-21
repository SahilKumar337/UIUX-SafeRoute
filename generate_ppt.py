"""
SafeRoute HCD Presentation Generator — v4 Dark Pro
Generates a 15-slide dark-themed PPTX matching the actual prototype UI.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt

# Try to use freeform; fall back gracefully
try:
    from pptx.util import Emu
except ImportError:
    pass

# ─────────────────────────────────────────────────────
#  DESIGN TOKENS  (match the React prototype exactly)
# ─────────────────────────────────────────────────────
BG0   = RGBColor(0x0B, 0x0E, 0x14)   # deepest bg
BG1   = RGBColor(0x13, 0x17, 0x20)
BG2   = RGBColor(0x1C, 0x21, 0x30)
BG3   = RGBColor(0x25, 0x2D, 0x3A)
BDR   = RGBColor(0x2A, 0x33, 0x47)   # border
GRN   = RGBColor(0x00, 0xD2, 0x6A)   # safety green
RED   = RGBColor(0xFF, 0x3D, 0x5A)   # SOS red
AMB   = RGBColor(0xFF, 0xC5, 0x42)   # hazard amber
BLU   = RGBColor(0x5B, 0x8D, 0xEF)   # info blue
PUR   = RGBColor(0x8B, 0x5C, 0xF6)   # accent purple
WHT   = RGBColor(0xFF, 0xFF, 0xFF)
T0    = RGBColor(0xFF, 0xFF, 0xFF)   # text primary
T1    = RGBColor(0xC6, 0xCE, 0xDF)   # text secondary
T2    = RGBColor(0x6B, 0x7A, 0x99)   # text muted
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────
#  HELPER UTILITIES
# ─────────────────────────────────────────────────────
def solid_bg(slide, color=BG0):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=1.0, radius=True):
    shape_type = 5  # MSO_SHAPE.ROUNDED_RECTANGLE = 5
    shp = slide.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(border_pt)
    else:
        shp.line.fill.background()
    return shp

def dim_color(col):
    """Return a very dim (~14% opacity) version of an RGBColor for badge backgrounds."""
    h = str(col)   # e.g. '00D26A'
    r = max(int(int(h[0:2], 16) * 0.14), 0x0B)
    g = max(int(int(h[2:4], 16) * 0.14), 0x0E)
    b = max(int(int(h[4:6], 16) * 0.14), 0x14)
    return RGBColor(r, g, b)

def feat_row_clean(sl, x, y, w, col, icon, title, body, divider=True):
    """
    Web-style FeatureRow: small colored icon badge (left) + bold title + gray description.
    Matches the React <FeatureRow> component exactly.
    """
    ROW_H    = 0.72   # total height per row
    BADGE_S  = 0.38   # badge square size
    BADGE_Y  = y + (ROW_H - BADGE_S) / 2  # vertically centered in row
    TEXT_X   = x + BADGE_S + 0.14
    TEXT_W   = w - BADGE_S - 0.14

    # Icon badge — dim colored square with colored icon
    badge = slide_shapes_add_rect(sl, x, BADGE_Y, BADGE_S, BADGE_S, dim_color(col))
    badge.text_frame.margin_top = Inches(0.03)
    badge.text_frame.margin_left = Inches(0.0)
    badge_p = badge.text_frame.paragraphs[0]
    badge_p.text = icon
    badge_p.font.size = Pt(15)
    badge_p.alignment = PP_ALIGN.CENTER

    # Title
    ttb = add_tb(sl, TEXT_X, y + 0.08, TEXT_W, 0.26)
    ttf = ttb.text_frame
    ttf.margin_left = ttf.margin_top = Inches(0.0)
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.name = 'Inter'; tp.font.size = Pt(12.5); tp.font.bold = True
    tp.font.color.rgb = T0

    # Description
    btb = add_tb(sl, TEXT_X, y + 0.37, TEXT_W, 0.28)
    btf = btb.text_frame
    btf.margin_left = btf.margin_top = Inches(0.0)
    bp = btf.paragraphs[0]
    bp.text = body
    bp.font.name = 'Inter'; bp.font.size = Pt(9.5)
    bp.font.color.rgb = T2

    # Thin horizontal divider line at bottom of row
    if divider:
        div = sl.shapes.add_shape(1, Inches(x), Inches(y + ROW_H - 0.01),
                                  Inches(w), Inches(0.008))
        div.fill.solid()
        div.fill.fore_color.rgb = BDR
        div.line.fill.background()

def slide_shapes_add_rect(sl, l, t, w, h, fill_col):
    """Straight (non-rounded) small rectangle helper used in feat_row_clean."""
    shp = sl.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_col
    shp.line.fill.background()
    return shp



def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def para(tf, text, size=12, bold=False, color=T1, italic=False,
         align=PP_ALIGN.LEFT, space_after=0, font='Inter', add=True):
    p = tf.add_paragraph() if add else tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    p.line_spacing = 1.2
    return p

def colored_card(slide, l, t, w, h, title, title_color, bullets,
                 bg=BG2, border=None, bullet_color=T1, size=11):
    """Draw a card with a colored title and bullet list."""
    shp = add_rect(slide, l, t, w, h, bg, border or BDR, 1.0)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_bottom = Inches(0.12)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Inter'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_after = Pt(7)
    for i, b in enumerate(bullets):
        bp = tf.add_paragraph()
        bp.text = f"•  {b}"
        bp.font.name = 'Inter'
        bp.font.size = Pt(size)
        bp.font.color.rgb = bullet_color
        bp.space_after = Pt(5)
        bp.line_spacing = 1.2
    return shp

def phase_tag(slide, phase, color, l=0.5, t=0.25):
    tb = add_tb(slide, l, t, 12.0, 0.35)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"  {phase.upper()}  "
    p.font.name = 'Inter'
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    p.line_spacing = 1.0

def slide_header(slide, phase, phase_color, title, subtitle):
    """Colored left straight vertical accent bar + title block."""
    # Vertical straight line (shape type 1 = RECTANGLE)
    bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.44), Inches(0.04), Inches(0.92))
    bar.fill.solid()
    bar.fill.fore_color.rgb = phase_color
    bar.line.fill.background()

    # Title and subtitle textbox
    tb = add_tb(slide, 0.72, 0.38, 12.0, 1.1)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = Inches(0.0)
    
    p1 = tf.paragraphs[0]
    p1.text = phase.upper()
    p1.font.name = 'Inter'; p1.font.size = Pt(8.5); p1.font.bold = True
    p1.font.color.rgb = phase_color; p1.space_after = Pt(2)
    
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.name = 'Inter'; p2.font.size = Pt(23); p2.font.bold = True
    p2.font.color.rgb = T0; p2.space_after = Pt(2)
    
    p3 = tf.add_paragraph()
    p3.text = subtitle
    p3.font.name = 'Inter'; p3.font.size = Pt(11)
    p3.font.color.rgb = T2

def footer(slide, cur, total):
    tb = add_tb(slide, 0.55, 7.12, 12.23, 0.25)
    tf = tb.text_frame
    tf.margin_left = tf.margin_top = 0
    # Thin straight line above footer (shape type 1 = RECTANGLE)
    ln = slide.shapes.add_shape(1, Inches(0.55), Inches(7.06), Inches(12.23), Inches(0.01))
    ln.fill.solid()
    ln.fill.fore_color.rgb = BDR
    ln.line.fill.background()
    
    p = tf.paragraphs[0]
    p.text = f"SafeRoute HCD Case Study  ·  Slide {cur} of {total}  ·  Design Jury 2026"
    p.font.name = 'Inter'; p.font.size = Pt(8); p.font.color.rgb = T2

def stat_box(slide, l, t, w, h, value, label, color):
    shp = add_rect(slide, l, t, w, h, BG2, color, 1.2)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.10)
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.name = 'Inter'; p1.font.size = Pt(34); p1.font.bold = True
    p1.font.color.rgb = color; p1.alignment = PP_ALIGN.CENTER; p1.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = 'Inter'; p2.font.size = Pt(9.5)
    p2.font.color.rgb = T2; p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.2

def draw_phone_mockup(slide, l, t, screen_type='dashboard'):
    """
    Draw a mini phone placeholder frame for screenshot insertion.
    l, t = top-left position in Inches
    """
    PW, PH = 2.3, 4.7   # phone outer size
    
    # Phone outer frame
    bezel = add_rect(slide, l, t, PW, PH, BG1, BDR, 2.0)
    
    # Dashed inner border card
    inner_card = slide.shapes.add_shape(5, Inches(l+0.06), Inches(t+0.06), Inches(PW-0.12), Inches(PH-0.12))
    inner_card.fill.solid()
    inner_card.fill.fore_color.rgb = BG0
    inner_card.line.color.rgb = AMB
    inner_card.line.width = Pt(1.5)
    inner_card.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # Centered text box
    tb = add_tb(slide, l+0.1, t+PH/2 - 0.7, PW-0.2, 1.4)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_top = Inches(0.02)
    
    screen_names = {
        'dashboard': 'Route Selection Screen',
        'navigation': 'Active Navigation Screen',
        'sos': 'SOS Emergency Screen',
        'hazard': 'Hazard Report Screen'
    }
    screen_name = screen_names.get(screen_type, 'Prototype Screen')
    
    para(tf, 'PLACEHOLDER', size=8, bold=True, color=AMB, align=PP_ALIGN.CENTER, add=False)
    para(tf, screen_name, size=10, bold=True, color=T0, align=PP_ALIGN.CENTER, add=True)
    para(tf, 'Insert screenshot here', size=8.5, color=T2, align=PP_ALIGN.CENTER, add=True)


def _phone_dashboard(slide, l, t, PW, PH, PAD, IW, CONTENT_TOP):
    # App header
    hdr = add_rect(slide, l+PAD, CONTENT_TOP, IW, 0.28, BG0, BDR, 0.5)
    tb = add_tb(slide, l+PAD+0.07, CONTENT_TOP+0.04, 1.0, 0.2)
    para(tb.text_frame, '🛡 SafeRoute', size=8, bold=True, color=T0, add=False)
    # avatar dot
    av = add_rect(slide, l+PW-PAD-0.28, CONTENT_TOP+0.04, 0.2, 0.2, PUR)

    # Dest bar
    DEST_TOP = CONTENT_TOP + 0.28
    dest = add_rect(slide, l+PAD, DEST_TOP, IW, 0.22, BG1, BDR, 0.5)
    dot = add_rect(slide, l+PAD+0.06, DEST_TOP+0.07, 0.08, 0.08, GRN)
    tb2 = add_tb(slide, l+PAD+0.20, DEST_TOP+0.04, 1.6, 0.16)
    para(tb2.text_frame, 'Campus Apartment', size=7, bold=True, color=T0, add=False)

    # Map area
    MAP_TOP = DEST_TOP + 0.22
    MAP_H = 1.75
    map_bg = add_rect(slide, l+PAD, MAP_TOP, IW, MAP_H, RGBColor(0x0D,0x11,0x17), BDR, 0.5)

    # Grid lines
    for gi in range(1, 4):
        lh = add_rect(slide, l+PAD, MAP_TOP + gi*(MAP_H/4), IW, 0.01, BDR)
        lv = add_rect(slide, l+PAD + gi*(IW/4), MAP_TOP, 0.01, MAP_H, BDR)

    # Safe route glow (drawn first underneath)
    g1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.45), Inches(l+PAD+0.37), Inches(MAP_TOP+0.30))
    g1.line.color.rgb = RGBColor(0x00, 0x50, 0x28)
    g1.line.width = Pt(6)
    
    g2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+0.30), Inches(l+PAD+1.00), Inches(MAP_TOP+0.30))
    g2.line.color.rgb = RGBColor(0x00, 0x50, 0x28)
    g2.line.width = Pt(6)

    # Safe route main line (green L-shape)
    s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.45), Inches(l+PAD+0.37), Inches(MAP_TOP+0.30))
    s1.line.color.rgb = GRN
    s1.line.width = Pt(2.5)

    s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+0.30), Inches(l+PAD+1.00), Inches(MAP_TOP+0.30))
    s2.line.color.rgb = GRN
    s2.line.width = Pt(2.5)

    # Unsafe route (red diagonal dashed line)
    u1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.45), Inches(l+PAD+1.00), Inches(MAP_TOP+0.30))
    u1.line.color.rgb = RED
    u1.line.width = Pt(1.5)
    u1.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # Dest pin label
    pin_tb = add_tb(slide, l+PAD+0.88, MAP_TOP+0.18, 1.0, 0.18)
    para(pin_tb.text_frame, '📍 Campus', size=6, bold=True, color=T0, add=False)
    pin_dot = add_rect(slide, l+PAD+0.98, MAP_TOP+0.30, 0.08, 0.08, GRN)

    # You are here
    you_tb = add_tb(slide, l+PAD+0.20, MAP_TOP+1.35, 0.9, 0.18)
    para(you_tb.text_frame, 'You are here', size=6, bold=True, color=PUR, add=False)
    you_dot = add_rect(slide, l+PAD+0.34, MAP_TOP+1.47, 0.10, 0.10, PUR)

    # Alley warning
    alley_tb = add_tb(slide, l+PAD+0.50, MAP_TOP+0.62, 0.8, 0.16)
    para(alley_tb.text_frame, '⚠ Dim Alley', size=6, bold=True, color=AMB, add=False)

    # Legend overlay
    leg_top = MAP_TOP + MAP_H - 0.38
    leg1 = add_rect(slide, l+PAD+0.05, leg_top, 0.85, 0.16, RGBColor(0x06,0x09,0x10), GRN, 0.5)
    ltb1 = add_tb(slide, l+PAD+0.08, leg_top+0.02, 0.80, 0.14)
    para(ltb1.text_frame, '— Safe  2.8km 34min', size=6, color=GRN, add=False)
    leg2 = add_rect(slide, l+PAD+0.05, leg_top+0.18, 0.85, 0.16, RGBColor(0x06,0x09,0x10), RED, 0.5)
    ltb2 = add_tb(slide, l+PAD+0.08, leg_top+0.20, 0.80, 0.14)
    para(ltb2.text_frame, '-- Unsafe 2.2km 27min', size=6, color=RED, add=False)

    # Route cards
    CARDS_TOP = MAP_TOP + MAP_H + 0.05
    # Safe card
    safe_card = add_rect(slide, l+PAD, CARDS_TOP, IW, 0.30, RGBColor(0x00,0x28,0x18), GRN, 0.8)
    stb = add_tb(slide, l+PAD+0.07, CARDS_TOP+0.04, 1.2, 0.22)
    para(stb.text_frame, 'SafeRoute  94% Safe', size=7, bold=True, color=GRN, add=False)
    stb2 = add_tb(slide, l+PAD+0.07, CARDS_TOP+0.16, 1.0, 0.14)
    para(stb2.text_frame, 'Lit streets · CCTV', size=6, color=T2, add=False)
    stb3 = add_tb(slide, l+PW-PAD-0.45, CARDS_TOP+0.04, 0.4, 0.24)
    para(stb3.text_frame, '34m', size=9, bold=True, color=T0, add=False)

    # Unsafe card
    un_card = add_rect(slide, l+PAD, CARDS_TOP+0.33, IW, 0.28, BG1, BDR, 0.5)
    utb = add_tb(slide, l+PAD+0.07, CARDS_TOP+0.37, 1.2, 0.20)
    para(utb.text_frame, 'Shortest  38%', size=7, bold=True, color=RED, add=False)
    utb2 = add_tb(slide, l+PAD+0.07, CARDS_TOP+0.49, 1.0, 0.12)
    para(utb2.text_frame, '⚠ Dim alleys', size=6, color=T2, add=False)
    utb3 = add_tb(slide, l+PW-PAD-0.45, CARDS_TOP+0.37, 0.4, 0.22)
    para(utb3.text_frame, '27m', size=9, bold=True, color=T0, add=False)

    # CTA button
    CTA_TOP = CARDS_TOP + 0.65
    cta = add_rect(slide, l+PAD, CTA_TOP, IW, 0.28, GRN)
    ctb = add_tb(slide, l+PAD+0.1, CTA_TOP+0.05, IW-0.2, 0.18)
    para(ctb.text_frame, '🛡 Start Safe Navigation', size=8, bold=True, color=BG0,
         align=PP_ALIGN.CENTER, add=False)


def _phone_navigation(slide, l, t, PW, PH, PAD, IW, CONTENT_TOP):
    # HUD
    hud = add_rect(slide, l+PAD, CONTENT_TOP, IW, 0.35, BG0, BDR, 0.5)
    # Back arrow
    bk = add_rect(slide, l+PAD+0.06, CONTENT_TOP+0.07, 0.18, 0.18, BG2)
    bktb = add_tb(slide, l+PAD+0.08, CONTENT_TOP+0.08, 0.16, 0.14)
    para(bktb.text_frame, '‹', size=10, color=T1, add=False)
    # Turn info
    t1 = add_tb(slide, l+PAD+0.28, CONTENT_TOP+0.04, 1.2, 0.14)
    para(t1.text_frame, 'SAFE ROUTE · 2.8km', size=6, bold=True, color=GRN, add=False)
    t2 = add_tb(slide, l+PAD+0.28, CONTENT_TOP+0.18, 1.2, 0.16)
    para(t2.text_frame, 'Turn right → Cedar Ave', size=7, bold=True, color=T0, add=False)
    # ETA box
    eta = add_rect(slide, l+PW-PAD-0.38, CONTENT_TOP+0.06, 0.32, 0.24, BG2)
    etatb = add_tb(slide, l+PW-PAD-0.36, CONTENT_TOP+0.06, 0.30, 0.24)
    para(etatb.text_frame, '34\nmin', size=7, bold=True, color=T0, align=PP_ALIGN.CENTER, add=False)

    # Map
    MAP_TOP = CONTENT_TOP + 0.35
    MAP_H = 2.65
    map_bg = add_rect(slide, l+PAD, MAP_TOP, IW, MAP_H, RGBColor(0x0D,0x11,0x17), BDR, 0.5)
    for gi in range(1, 4):
        add_rect(slide, l+PAD, MAP_TOP + gi*(MAP_H/4), IW, 0.01, BDR)
        add_rect(slide, l+PAD + gi*(IW/4), MAP_TOP, 0.01, MAP_H, BDR)

    # Safe route glow (thick line)
    g1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.85), Inches(l+PAD+0.37), Inches(MAP_TOP+0.30))
    g1.line.color.rgb = RGBColor(0x00, 0x40, 0x20)
    g1.line.width = Pt(8)
    
    g2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+0.30), Inches(l+PAD+1.10), Inches(MAP_TOP+0.30))
    g2.line.color.rgb = RGBColor(0x00, 0x40, 0x20)
    g2.line.width = Pt(8)

    # Safe route main line (green L-shape)
    s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.85), Inches(l+PAD+0.37), Inches(MAP_TOP+0.30))
    s1.line.color.rgb = GRN
    s1.line.width = Pt(3.5)

    s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+0.30), Inches(l+PAD+1.10), Inches(MAP_TOP+0.30))
    s2.line.color.rgb = GRN
    s2.line.width = Pt(3.5)

    # Dim unsafe route (dashed)
    u1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.85), Inches(l+PAD+1.10), Inches(MAP_TOP+0.30))
    u1.line.color.rgb = RGBColor(0x80, 0x20, 0x30)
    u1.line.width = Pt(1.5)
    u1.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # Pins
    ptb = add_tb(slide, l+PAD+0.90, MAP_TOP+0.10, 0.9, 0.16)
    para(ptb.text_frame, '📍 Campus', size=6, bold=True, color=T0, add=False)
    add_rect(slide, l+PAD+1.00, MAP_TOP+0.22, 0.08, 0.08, GRN)
    youhd = add_tb(slide, l+PAD+0.15, MAP_TOP+1.72, 0.9, 0.16)
    para(youhd.text_frame, 'You are here', size=6, color=PUR, add=False)
    add_rect(slide, l+PAD+0.30, MAP_TOP+1.84, 0.10, 0.10, PUR)
    # Safety badge
    sb = add_rect(slide, l+PAD+0.06, MAP_TOP+MAP_H-0.24, 1.0, 0.18, RGBColor(0x06,0x09,0x10), GRN, 0.5)
    sbtb = add_tb(slide, l+PAD+0.09, MAP_TOP+MAP_H-0.22, 0.94, 0.14)
    para(sbtb.text_frame, '● 94% Safety Corridor', size=6, color=GRN, add=False)

    # Bottom controls
    CTRL_TOP = MAP_TOP + MAP_H + 0.05
    ctrl_bg = add_rect(slide, l+PAD, CTRL_TOP, IW, 0.42, BG0, BDR, 0.5)
    # Report
    rpt = add_rect(slide, l+PAD+0.05, CTRL_TOP+0.08, 0.48, 0.26, BG2, BDR, 0.5)
    rtb = add_tb(slide, l+PAD+0.08, CTRL_TOP+0.11, 0.42, 0.18)
    para(rtb.text_frame, '⚠ Report', size=6, bold=True, color=AMB, add=False)
    # SOS pulsing circle
    sos = add_rect(slide, l+PAD+0.68, CTRL_TOP+0.04, 0.36, 0.34, RED)
    sostb = add_tb(slide, l+PAD+0.70, CTRL_TOP+0.10, 0.32, 0.18)
    para(sostb.text_frame, 'SOS', size=7, bold=True, color=WHT, align=PP_ALIGN.CENTER, add=False)
    # Share
    shr = add_rect(slide, l+PW-PAD-0.56, CTRL_TOP+0.08, 0.48, 0.26, BG2, BDR, 0.5)
    shtb = add_tb(slide, l+PW-PAD-0.53, CTRL_TOP+0.11, 0.42, 0.18)
    para(shtb.text_frame, '↗ Share', size=6, bold=True, color=BLU, add=False)


def _phone_sos(slide, l, t, PW, PH, PAD, IW, CONTENT_TOP):
    # Subtle red tinted bg overlay
    tint = add_rect(slide, l+PAD, CONTENT_TOP, IW, PH-PAD*2-0.12, RGBColor(0x18,0x08,0x0C))
    # Header
    add_rect(slide, l+PAD, CONTENT_TOP, IW, 0.26, BG0, BDR, 0.5)
    htb = add_tb(slide, l+PAD+0.10, CONTENT_TOP+0.05, 1.5, 0.18)
    para(htb.text_frame, '🚨 SOS Alert', size=8, bold=True, color=RED, add=False)

    # Countdown circle
    CC_TOP = CONTENT_TOP + 0.38
    cc = add_rect(slide, l+PAD+0.45, CC_TOP, 1.05, 1.05, RGBColor(0x28,0x08,0x10), RED, 2.0)
    cctb = add_tb(slide, l+PAD+0.55, CC_TOP+0.12, 0.85, 0.65)
    para(cctb.text_frame, '3', size=40, bold=True, color=RED, align=PP_ALIGN.CENTER, add=False)

    # Sending text
    st = add_tb(slide, l+PAD+0.12, CC_TOP+1.15, IW-0.24, 0.20)
    para(st.text_frame, 'Sending SOS Alert...', size=8, bold=True, color=T0, align=PP_ALIGN.CENTER, add=False)
    st2 = add_tb(slide, l+PAD+0.12, CC_TOP+1.36, IW-0.24, 0.24)
    para(st2.text_frame, 'GPS sent to contacts in 3 seconds', size=7, color=T2, align=PP_ALIGN.CENTER, add=False)

    # Checklist
    CL_TOP = CC_TOP + 1.68
    for i, (label, done) in enumerate([('GPS lock acquired', True),
                                        ('Composing SMS', False),
                                        ('Alerting campus security', False)]):
        dot_color = GRN if done else BG3
        dot = add_rect(slide, l+PAD+0.12, CL_TOP+i*0.22, 0.14, 0.14, dot_color)
        if done:
            dottb = add_tb(slide, l+PAD+0.13, CL_TOP+i*0.22+0.01, 0.12, 0.12)
            para(dottb.text_frame, '✓', size=6, bold=True, color=BG0, add=False)
        ltb = add_tb(slide, l+PAD+0.30, CL_TOP+i*0.22, 1.2, 0.18)
        para(ltb.text_frame, label, size=6.5, color=T0 if done else T2, add=False)

    # Abort button
    AB_TOP = CL_TOP + 0.75
    ab = add_rect(slide, l+PAD+0.12, AB_TOP, IW-0.24, 0.26, BG0, RED, 1.0)
    abtb = add_tb(slide, l+PAD+0.15, AB_TOP+0.06, IW-0.30, 0.16)
    para(abtb.text_frame, 'Hold to Cancel Alert', size=7, bold=True, color=RED, align=PP_ALIGN.CENTER, add=False)


def _phone_hazard(slide, l, t, PW, PH, PAD, IW, CONTENT_TOP):
    # Status + small map
    add_rect(slide, l+PAD, CONTENT_TOP, IW, 0.26, BG0, BDR, 0.5)
    htb = add_tb(slide, l+PAD+0.08, CONTENT_TOP+0.05, 1.5, 0.18)
    para(htb.text_frame, 'Navigation', size=8, bold=True, color=T0, add=False)

    MAP_TOP = CONTENT_TOP + 0.26
    MAP_H = 1.60
    add_rect(slide, l+PAD, MAP_TOP, IW, MAP_H, RGBColor(0x0D,0x11,0x17), BDR, 0.5)
    for gi in range(1, 4):
        add_rect(slide, l+PAD, MAP_TOP + gi*(MAP_H/4), IW, 0.01, BDR)
        add_rect(slide, l+PAD + gi*(IW/4), MAP_TOP, 0.01, MAP_H, BDR)
    # Dimmed safe route (green L-shape)
    s1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+1.30), Inches(l+PAD+0.37), Inches(MAP_TOP+0.30))
    s1.line.color.rgb = RGBColor(0x00, 0x60, 0x38)
    s1.line.width = Pt(2.0)

    s2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(l+PAD+0.37), Inches(MAP_TOP+0.30), Inches(l+PAD+1.00), Inches(MAP_TOP+0.30))
    s2.line.color.rgb = RGBColor(0x00, 0x60, 0x38)
    s2.line.width = Pt(2.0)
    # Hazard pins
    h1 = add_tb(slide, l+PAD+0.50, MAP_TOP+0.30, 1.0, 0.16)
    para(h1.text_frame, '⚠ Dim Lights', size=6, bold=True, color=AMB, add=False)
    h2 = add_tb(slide, l+PAD+0.20, MAP_TOP+0.75, 1.0, 0.16)
    para(h2.text_frame, '⚠ Blocked Path', size=6, bold=True, color=AMB, add=False)

    # Bottom sheet modal
    MODAL_TOP = MAP_TOP + MAP_H + 0.05
    modal = add_rect(slide, l+PAD, MODAL_TOP, IW, 2.10, BG1, BDR, 1.0)
    # Handle
    handle = add_rect(slide, l+PW/2-0.20, MODAL_TOP+0.08, 0.40, 0.04, BDR)
    # Title
    mtb = add_tb(slide, l+PAD+0.10, MODAL_TOP+0.18, IW-0.20, 0.20)
    para(mtb.text_frame, 'Report Safety Hazard', size=9, bold=True, color=T0, add=False)
    # 4 category chips (2x2 grid)
    CATS = ['Dim Lighting', 'Blocked Path', 'Suspicious', 'Unsafe Road']
    for ci, cat in enumerate(CATS):
        cx = l+PAD+0.08 + (ci%2)*(IW/2-0.06)
        cy = MODAL_TOP + 0.44 + (ci//2)*0.30
        chip = add_rect(slide, cx, cy, IW/2-0.12, 0.24,
                        RGBColor(0x24,0x14,0x48) if ci==0 else BG2,
                        PUR if ci==0 else BDR, 0.7)
        ctb = add_tb(slide, cx+0.05, cy+0.05, IW/2-0.22, 0.14)
        para(ctb.text_frame, cat, size=6.5, bold=True,
             color=PUR if ci==0 else T2, align=PP_ALIGN.CENTER, add=False)
    # Submit button
    SUB_TOP = MODAL_TOP + 1.12
    sub = add_rect(slide, l+PAD+0.08, SUB_TOP, IW-0.16, 0.28, AMB)
    stb = add_tb(slide, l+PAD+0.12, SUB_TOP+0.06, IW-0.24, 0.18)
    para(stb.text_frame, '📍 Publish Hazard Report', size=7.5, bold=True, color=BG0,
         align=PP_ALIGN.CENTER, add=False)


# ─────────────────────────────────────────────────────
#  BUILD PRESENTATION
# ─────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    TOTAL = 15

    # ════════════════════════════════════════════════
    # SLIDE 1 — TITLE / HERO
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Introduction', PUR, 'SafeRoute', 'Smart Safety Navigation for Solo Night Travelers')

    # Intro description paragraph
    desc_tb = add_tb(sl, 0.68, 1.55, 7.3, 0.7)
    para(desc_tb.text_frame, 
         'A Human-Centered Design case study — a mobile navigation app that routes solo pedestrians '
         'along well-lit, high-footfall streets instead of the fastest dark shortcuts.', 
         size=11, color=T1, add=False)

    # Inline pills
    pills = [('Smart Cities', PUR), ('Social Safety', GRN), ('HCD 5-Phase', BLU), ('Live Prototype', AMB)]
    for pi, (label, col) in enumerate(pills):
        px = 0.68 + pi * 1.8
        badge = add_rect(sl, px, 2.26, 1.6, 0.28, BG2, col, 0.8)
        btb = add_tb(sl, px+0.05, 2.30, 1.5, 0.2)
        para(btb.text_frame, label, size=8, bold=True, color=col, align=PP_ALIGN.CENTER, add=False)

    # 2x2 grid of cards
    stats = [
        ('5', 'HCD Phases', PUR, 0.68, 2.80),
        ('73%', 'Women anxious walking alone', RED, 4.30, 2.80),
        ('94%', 'App safety score achieved', GRN, 0.68, 4.35),
        ('3s', 'SOS response trigger', AMB, 4.30, 4.35),
    ]
    for v, lb, col, sx, sy in stats:
        stat_box(sl, sx, sy, 3.4, 1.35, v, lb, col)

    # Right side phone mockup
    draw_phone_mockup(sl, 8.8, 1.4, 'dashboard')
    cap_tb = add_tb(sl, 8.0, 6.2, 3.9, 0.5)
    para(cap_tb.text_frame, 'Route Selection\nSafe vs unsafe routes with real distances', size=9, color=T2, align=PP_ALIGN.CENTER, add=False)

    footer(sl, 1, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 2 — DISCOVER: The Problem
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 1 – Discover', BLU, 'The Darkness Trap', 'Problem Identification & Secondary Research')

    add_rect(sl, 0.5, 1.7, 12.33, 0.02, BDR)   # divider

    # Intro text
    tb = add_tb(sl, 0.55, 1.82, 12.0, 0.45)
    para(tb.text_frame,
         'Traditional mapping apps optimize for SPEED ONLY. At night, pedestrians are regularly routed through unlit alleys, dark parks, '
         'and back lanes — not because it is safe, but because it saves 2 minutes.',
         size=12, color=T1, add=False)

    # 3 stat boxes
    for si, (v, lb, col) in enumerate([
        ('73%', 'Of women feel highly anxious\nwalking home alone at night', RED),
        ('85%', 'Prefer well-lit paths even if\n5–10 min longer', AMB),
        ('91%', 'Define safety by active\nstorefronts & streetlights', GRN),
    ]):
        stat_box(sl, 0.55 + si * 4.1, 2.45, 3.8, 1.65, v, lb, col)

    # Root cause box
    colored_card(sl, 0.55, 4.25, 12.0, 2.6, 'Root Cause: The Speed Bias', AMB,
                 ['Google Maps calculates routes by fastest time — lighting data is NEVER considered.',
                  'There is no "safe walk" mode for pedestrians in any major navigation app today.',
                  'Community safety data (incidents, dim lights, blocked paths) is not integrated into live route logic.',
                  'The result: Elena gets routed through a pitch-black park to save 90 seconds.'],
                 bg=BG2, border=AMB)
    footer(sl, 2, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 3 — DISCOVER: Competitor Analysis
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 1 – Discover', BLU, 'Competitor Analysis', 'What Existing Apps Fail to Solve')

    intro = add_tb(sl, 0.55, 1.75, 12.0, 0.35)
    para(intro.text_frame,
         'Existing tools offer location sharing but fail to provide preventative, safety-optimized routing.',
         size=12, color=T1, add=False)

    # Table
    headers = ['Feature', 'Google Maps', 'Life360', 'bSafe', 'SafeRoute  ✦']
    rows = [
        ['Safety Routing',     '✗ Speed only',       '✗ No routing',       '✗ SOS only',      '✓ Illumination-based'],
        ['Dark Area Warnings', '✗ None',              '✗ None',             '✗ None',           '✓ Real-time hazard flags'],
        ['One-Tap SOS',        '✗ Open dialer',       '✓ Basic notif',      '✓ Basic alarm',    '✓ Siren + SMS + GPS'],
        ['Community Reports',  '✓ Traffic only',      '✗ None',             '✗ None',           '✓ Street-level hazards'],
        ['Night-Opt UI',       '✗ Standard',          '✗ Standard',         '✗ Standard',       '✓ Dark + Glow Map'],
    ]
    ROW_H = 0.60
    COL_WIDTHS = [2.4, 2.15, 2.15, 2.15, 2.15]
    TABLE_L = 0.55; TABLE_T = 2.18

    # Header row
    x = TABLE_L
    for ci, (h, cw) in enumerate(zip(headers, COL_WIDTHS)):
        hdr = add_rect(sl, x, TABLE_T, cw, ROW_H * 0.70, BG2 if ci < 4 else RGBColor(0x00,0x28,0x18), GRN if ci == 4 else BDR, 0.8)
        htb = add_tb(sl, x+0.07, TABLE_T+0.10, cw-0.14, 0.45)
        para(htb.text_frame, h, size=10, bold=True, color=GRN if ci == 4 else T1,
             align=PP_ALIGN.CENTER, add=False)
        x += cw + 0.02

    # Data rows
    for ri, row in enumerate(rows):
        x = TABLE_L
        for ci, (cell, cw) in enumerate(zip(row, COL_WIDTHS)):
            is_safe = ci == 4
            good = '✓' in cell; bad = '✗' in cell
            bg = RGBColor(0x00,0x1A,0x0F) if (is_safe and good) else (BG2 if ri%2==0 else BG1)
            bdr = GRN if is_safe else BDR
            card = add_rect(sl, x, TABLE_T + (ri+1)*ROW_H*0.70 + ri*0.04, cw, ROW_H*0.62, bg, bdr, 0.5)
            ctb = add_tb(sl, x+0.07, TABLE_T + (ri+1)*ROW_H*0.70 + ri*0.04 + 0.10, cw-0.14, 0.45)
            col = GRN if good else (RED if bad else T2)
            if ci == 0: col = T1
            if is_safe: col = GRN if good else RED
            para(ctb.text_frame, cell, size=9.5, color=col, bold=is_safe and good,
                 align=PP_ALIGN.CENTER, add=False)
            x += cw + 0.02

    footer(sl, 3, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 4 — DEFINE: User Persona
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 2 – Define', PUR, 'Meet Elena Rivera', 'Primary User Persona — The Solo Student')

    # Avatar column
    av = add_rect(sl, 0.55, 1.75, 2.5, 5.15, BG2, PUR, 1.0)
    avtf = av.text_frame; avtf.word_wrap = True
    avtf.margin_top = Inches(0.35)
    p = avtf.paragraphs[0]
    p.text = '👩‍🎓'; p.font.size = Pt(56); p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(8)
    p2 = avtf.add_paragraph()
    p2.text = 'Elena Rivera'; p2.font.name='Inter'; p2.font.size=Pt(16); p2.font.bold=True
    p2.font.color.rgb = T0; p2.alignment=PP_ALIGN.CENTER; p2.space_after=Pt(3)
    p3 = avtf.add_paragraph()
    p3.text = 'Age 21 · College Student'; p3.font.name='Inter'; p3.font.size=Pt(10)
    p3.font.color.rgb = T2; p3.alignment=PP_ALIGN.CENTER; p3.space_after=Pt(14)
    p4 = avtf.add_paragraph()
    p4.text = '"Walking back at 10 PM is terrifying. I clutch my keys and hope for the best."'
    p4.font.name='Inter'; p4.font.size=Pt(9); p4.font.italic=True
    p4.font.color.rgb = T1; p4.alignment=PP_ALIGN.CENTER; p4.line_spacing=1.3

    # Right column cards
    RL = 3.3
    RW = 9.6
    colored_card(sl, RL, 1.75, RW, 1.30, 'Core Quote', PUR,
                 ['"Walking back to my dorm from my bookstore shift at 10 PM is always terrifying. '
                  'I clutch my keys and hope for the best. Google Maps routes me through the dark park '
                  'and doesn\'t care."'],
                 bg=RGBColor(0x18,0x0F,0x2A), border=PUR, size=11)
    colored_card(sl, RL, 3.18, RW/3-0.08, 1.60, 'Goals', GRN,
                 ['Avoid unlit shortcuts', 'Alert family fast', 'Skip expensive Ubers', 'Walk confidently at night'],
                 bg=BG2, border=GRN, size=10)
    colored_card(sl, RL + RW/3 + 0.04, 3.18, RW/3-0.08, 1.60, 'Behaviors', BLU,
                 ['Holds keys as weapon', 'Keeps one earbud out', 'Texts roommate ETA', 'Calls friend while walking'],
                 bg=BG2, border=BLU, size=10)
    colored_card(sl, RL + 2*RW/3 + 0.08, 3.18, RW/3-0.08, 1.60, 'Pain Points', RED,
                 ['Apps route through dark parks', 'Rideshares too costly', 'No lighting layer on maps', 'Fear is exhausting'],
                 bg=BG2, border=RED, size=10)
    colored_card(sl, RL, 4.94, RW, 1.85, 'Why SafeRoute Matters for Elena', AMB,
                 ['SafeRoute directly resolves ALL three pain points: it shows lit routes on a live map, '
                  'provides a one-tap SOS to alert contacts, and is free to use — unlike rideshares.',
                  'The 94% safety score gives Elena the confidence she has never had on a navigation app before.'],
                 bg=BG2, border=AMB, size=11)
    footer(sl, 4, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 5 — DEFINE: Empathy Map
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 2 – Define', PUR, 'Empathy Map', "Mapping Elena's Emotional Reality")

    Q = [
        ('SAYS 💬', BLU,
         ['"Should I pay $15 for a 5-min Uber?"',
          '"Let me text my roommate I\'m leaving work."',
          '"Why are these streetlamps broken again?"']),
        ('THINKS 🧠', PUR,
         ['"Is someone following me right now?"',
          '"I feel so vulnerable alone out here."',
          '"I hate being scared of my own neighbourhood."']),
        ('DOES 🚶', AMB,
         ['Walks fast, grips phone tightly',
          'Turns around at every sudden noise',
          'Calls a friend to stay on the phone']),
        ('FEELS ❤️', RED,
         ['Highly anxious in unlit areas',
          'Financially burdened by rideshares',
          'Hyper-vigilant of every movement']),
    ]
    positions = [(0.55, 1.75), (6.62, 1.75), (0.55, 4.35), (6.62, 4.35)]
    for (title, color, items), (lpos, tpos) in zip(Q, positions):
        colored_card(sl, lpos, tpos, 5.85, 2.35, title, color, items, bg=BG2, border=color, size=11)

    # Design implications strip
    add_rect(sl, 0.55, 6.80, 12.33, 0.04, PUR)
    impl = add_tb(sl, 0.55, 6.88, 12.0, 0.30)
    para(impl.text_frame,
         '→ Design Implications:  Illumination map layer  ·  3-second SOS panic button  ·  One-tap location share',
         size=10, bold=True, color=AMB, add=False)
    footer(sl, 5, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 6 — DEFINE: User Journey Map
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 2 – Define', PUR, 'The Walk Home', "Elena's User Journey — From Fear to Safety")

    STEPS = [
        ('1', 'Shift Ends', '10:00 PM. Anxious.\nShares location to roommate.', '😐', T2, 25),
        ('2', 'Open Maps', 'Standard app routes\nthrough dim park.', '😰', AMB, 55),
        ('3', 'Dark Alley', 'Reaches pitch-black\nshortcut. Peak fear.', '😨', RED, 90),
        ('4', 'SafeRoute', 'Switches app. 94%\nsafe route selected.', '🙂', GRN, 35),
        ('5', 'Home Safe', 'Arrives safely.\nConfidence restored.', '🤩', GRN, 10),
    ]
    SW = 2.35
    timeline_line = sl.shapes.add_shape(1, Inches(0.55), Inches(3.65), Inches(12.23), Inches(0.02))
    timeline_line.fill.solid()
    timeline_line.fill.fore_color.rgb = BDR
    timeline_line.line.fill.background()

    for si, (num, title, desc, emoji, col, anxiety) in enumerate(STEPS):
        sx = 0.55 + si * (SW + 0.1)
        card = add_rect(sl, sx, 1.75, SW, 3.8, BG2, col, 1.2)
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.20)
        p1 = tf.paragraphs[0]
        p1.text = num; p1.font.name='Inter'; p1.font.size=Pt(28); p1.font.bold=True
        p1.font.color.rgb = col; p1.alignment=PP_ALIGN.CENTER; p1.space_after=Pt(5)
        p2 = tf.add_paragraph()
        p2.text = title; p2.font.name='Inter'; p2.font.size=Pt(13); p2.font.bold=True
        p2.font.color.rgb = T0; p2.alignment=PP_ALIGN.CENTER; p2.space_after=Pt(6)
        p3 = tf.add_paragraph()
        p3.text = desc; p3.font.name='Inter'; p3.font.size=Pt(10)
        p3.font.color.rgb = T2; p3.alignment=PP_ALIGN.CENTER; p3.space_after=Pt(10); p3.line_spacing=1.3
        p4 = tf.add_paragraph()
        p4.text = emoji; p4.font.size=Pt(26); p4.alignment=PP_ALIGN.CENTER

        # Anxiety bar
        BAR_MAX_H = 1.35
        bar_h = anxiety / 100 * BAR_MAX_H
        bar_col = RED if anxiety > 70 else (AMB if anxiety > 40 else GRN)
        bar_top = 5.75 + BAR_MAX_H - bar_h
        add_rect(sl, sx + SW/2 - 0.3, bar_top, 0.6, bar_h, bar_col)
        bt = add_tb(sl, sx + SW/2 - 0.4, 7.12, 0.8, 0.18)
        para(bt.text_frame, title[:6], size=7, color=T2, align=PP_ALIGN.CENTER, add=False)

    lab = add_tb(sl, 0.55, 5.60, 1.0, 0.22)
    para(lab.text_frame, 'ANXIETY', size=7, color=T2, add=False)
    footer(sl, 6, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 7 — IDEATE: HMW & Crazy 8
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 3 – Ideate', AMB, 'Brainstorming Safety', '"How Might We" & Crazy 8 Concepts')

    colored_card(sl, 0.55, 1.75, 12.0, 1.50, 'How Might We (HMW) Statements', AMB,
                 ['HMW 1 (Routing): How might we guide solo nighttime walkers along lit streets without adding heavy travel delays?',
                  'HMW 2 (Emergency): How might we provide a frictionless, one-tap panic mechanism that alerts contacts and deters offenders in seconds?',
                  'HMW 3 (Community): How might we empower pedestrians to flag unsafe streets in real time, improving the map for everyone?'],
                 bg=RGBColor(0x1A,0x12,0x04), border=AMB, size=11)

    concepts = [
        ('🗺', 'Illumination HUD', 'Color-coded route lines: green (safe, lit streets) vs red dashed (dark alleys). Built into live map.', GRN),
        ('🚨', '3s SOS Countdown', 'Pulsing countdown + large abort button before GPS dispatch. Prevents false alarms.', RED),
        ('📍', 'Hazard Pin System', 'Tap to report dim lights, blocked paths, suspicious crowds. Pin appears on live map instantly.', AMB),
        ('🛡', 'Safety Score Badge', '"94% Safety Corridor" — derived from CCTV coverage, footfall & community reports.', PUR),
        ('📤', 'Live Location Share', 'One tap sends a live GPS link via SMS to pre-set emergency contacts.', BLU),
        ('🌙', 'Night-Mode Map', 'Dark CartoDB tiles + glowing neon routes. Preserves night vision while navigating dark streets.', T1),
    ]
    CW = 3.9; CGAP = 0.14
    for ci, (em, title, body, col) in enumerate(concepts):
        cx = 0.55 + (ci % 3) * (CW + CGAP)
        cy = 3.45 + (ci // 3) * 1.55
        card = add_rect(sl, cx, cy, CW, 1.40, BG2, col, 0.8)
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.12)
        p1 = tf.paragraphs[0]
        p1.text = f"{em}  {title}"; p1.font.name='Inter'; p1.font.size=Pt(11); p1.font.bold=True
        p1.font.color.rgb = col; p1.space_after=Pt(6)
        p2 = tf.add_paragraph()
        p2.text = body; p2.font.name='Inter'; p2.font.size=Pt(9.5)
        p2.font.color.rgb = T2; p2.line_spacing=1.2

    footer(sl, 7, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 8 — IDEATE: Information Architecture
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 3 – Ideate', AMB, 'App Architecture', 'Information Architecture & 4-Screen Design')

    IA_COLS = [
        ('1. Map HUD', GRN, ['Search destination', 'Safety score toggle', 'Route selector', 'Live hazard pins', 'Destination search']),
        ('2. Navigation', BLU, ['Turn-by-turn HUD', 'Safety corridor badge', 'SOS pulse button', 'Share location', 'Back to routes']),
        ('3. SOS Center', RED, ['3s countdown overlay', 'Emergency contacts', 'GPS SMS dispatch', 'Audio siren control', 'Abort button']),
        ('4. Hazard Reports', AMB, ['Category selection', 'Map pin placement', 'Notes text input', 'Community feed', 'Instant map update']),
    ]
    for ci, (title, col, items) in enumerate(IA_COLS):
        cx = 0.55 + ci * 3.1
        colored_card(sl, cx, 1.80, 2.85, 3.90, title, col, items, bg=BG2, border=col, size=10)

    # Task flow
    colored_card(sl, 0.55, 5.88, 12.0, 1.08, 'Primary Task Flow', GRN,
                 ['Open App  →  Set Destination  →  View Both Routes on Map  →  Compare Safety Score  →  '
                  'Select Safe Route  →  Start Navigation  →  SOS if Needed  →  Report Hazards  →  Arrive Safely'],
                 bg=BG2, border=GRN, size=11)
    footer(sl, 8, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 9 — DESIGN: Design System
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 4 – Design', GRN, 'Design System', 'Night-Optimized Visual Identity & Components')

    # Color tokens row
    tok_lab = add_tb(sl, 0.55, 1.75, 5.0, 0.22)
    para(tok_lab.text_frame, 'COLOR TOKENS', size=9, bold=True, color=T2, add=False)

    TOKENS = [
        (BG0, '#0B0E14', 'Stealth BG'), (GRN, '#00D26A', 'Safety Green'),
        (RED, '#FF3D5A', 'SOS Red'),    (AMB, '#FFC542', 'Hazard Amber'),
        (PUR, '#8B5CF6', 'Accent Purple'), (BLU, '#5B8DEF', 'Info Blue'),
    ]
    for ti, (col, hx, lb) in enumerate(TOKENS):
        tx = 0.55 + ti * 2.06
        swatch = add_rect(sl, tx, 2.00, 1.90, 0.65, col, WHT, 0.3)
        stb = add_tb(sl, tx, 2.70, 1.90, 0.35)
        para(stb.text_frame, f'{lb}\n{hx}', size=8, color=T2, align=PP_ALIGN.CENTER, add=False)

    # Typography
    typ_lab = add_tb(sl, 0.55, 3.22, 5.0, 0.22)
    para(typ_lab.text_frame, 'TYPOGRAPHY', size=9, bold=True, color=T2, add=False)
    typ_card = add_rect(sl, 0.55, 3.46, 7.0, 1.3, BG2, BDR, 0.8)
    ttf = typ_card.text_frame; ttf.word_wrap = True; ttf.margin_left = Inches(0.15); ttf.margin_top = Inches(0.10)
    p1 = ttf.paragraphs[0]; p1.text = 'SafeRoute'; p1.font.name='Inter'; p1.font.size=Pt(30)
    p1.font.bold=True; p1.font.color.rgb=T0; p1.space_after=Pt(4)
    p2 = ttf.add_paragraph(); p2.text = 'Inter Black — Headings'
    p2.font.name='Inter'; p2.font.size=Pt(9); p2.font.color.rgb=T2; p2.space_after=Pt(8)
    p3 = ttf.add_paragraph(); p3.text = 'Turn right in 200m → Cedar Ave'
    p3.font.name='Inter'; p3.font.size=Pt(13); p3.font.color.rgb=T1; p3.space_after=Pt(3)
    p4 = ttf.add_paragraph(); p4.text = 'Inter Regular — Body Text'
    p4.font.name='Inter'; p4.font.size=Pt(9); p4.font.color.rgb=T2

    # Principles
    prin_lab = add_tb(sl, 7.70, 3.22, 5.0, 0.22)
    para(prin_lab.text_frame, 'DESIGN PRINCIPLES', size=9, bold=True, color=T2, add=False)
    PRINS = [
        ('🌙 Night Vision First', 'Dark #0B0E14 BG preserves night vision — no blinding white screen while walking dark streets.'),
        ('⚡ Speed Over Beauty', 'All safety actions require ≤1 tap. No deep menus. Designed for dark, stressful conditions.'),
        ('🎯 Glanceable UI', 'Large fonts, glow effects, high contrast — readable at a single glance while moving.'),
    ]
    for pi, (ptitle, pbody) in enumerate(PRINS):
        pc = add_rect(sl, 7.70, 3.46 + pi * 0.95, 5.08, 0.85, BG2, BDR, 0.6)
        pctf = pc.text_frame; pctf.word_wrap = True
        pctf.margin_left = Inches(0.15); pctf.margin_top = Inches(0.10)
        pp1 = pctf.paragraphs[0]; pp1.text = ptitle
        pp1.font.name='Inter'; pp1.font.size=Pt(11); pp1.font.bold=True; pp1.font.color.rgb=T0; pp1.space_after=Pt(3)
        pp2 = pctf.add_paragraph(); pp2.text = pbody
        pp2.font.name='Inter'; pp2.font.size=Pt(9.5); pp2.font.color.rgb=T2; pp2.line_spacing=1.20

    # Component showcase
    comp_lab = add_tb(sl, 0.55, 4.88, 5.0, 0.22)
    para(comp_lab.text_frame, 'COMPONENT SHOWCASE', size=9, bold=True, color=T2, add=False)
    COMPS = [
        ('Primary CTA', GRN, BG0), ('Danger Button', RED, WHT),
        ('94% Safe', GRN, RGBColor(0,40,20)), ('38% Safe', RED, RGBColor(40,0,10)),
        ('⚠ Hazard', AMB, BG2), ('SOS', RED, WHT),
    ]
    for bi, (lbl, col, txt_col) in enumerate(COMPS):
        bx = 0.55 + bi * 2.06
        btn = add_rect(sl, bx, 5.14, 1.90, 0.38, col if bi < 2 else BG2, col, 1.0)
        btb = add_tb(sl, bx+0.08, 5.20, 1.74, 0.26)
        para(btb.text_frame, lbl, size=9, bold=True, color=BG0 if bi < 2 else col,
             align=PP_ALIGN.CENTER, add=False)

    footer(sl, 9, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 10 — PROTOTYPE: Screen 1 — Route Selection
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 4 – Prototype', GRN, 'Screen 1: Route Selection', 'The Core Decision — Safe vs Fastest Route')

    draw_phone_mockup(sl, 0.55, 1.50, 'dashboard')
    cap_tb = add_tb(sl, 0.20, 6.30, 3.0, 0.3)
    para(cap_tb.text_frame, 'Route Selection Screen', size=8, color=T2, align=PP_ALIGN.CENTER, add=False)

    # WHAT THIS SCREEN DOES — slim green banner
    RX = 3.25; RW = 9.50
    wtsd = add_rect(sl, RX, 1.50, RW, 0.90, RGBColor(0x00, 0x28, 0x18), GRN, 1.2)
    wtsdtf = wtsd.text_frame; wtsdtf.word_wrap = True
    wtsdtf.margin_left = Inches(0.15); wtsdtf.margin_top = Inches(0.09)
    wp = wtsdtf.paragraphs[0]; wp.text = '➤  WHAT THIS SCREEN DOES'
    wp.font.name = 'Inter'; wp.font.size = Pt(9.5); wp.font.bold = True; wp.font.color.rgb = GRN; wp.space_after = Pt(4)
    wp2 = wtsdtf.add_paragraph()
    wp2.text = 'Displays both routes simultaneously on a live, pannable dark map — letting the user instantly see WHERE the danger lies before choosing their path.'
    wp2.font.name = 'Inter'; wp2.font.size = Pt(10); wp2.font.color.rgb = T1; wp2.line_spacing = 1.30

    # Feature rows — web-style clean rows
    FEATS10 = [
        (GRN, '🗺', 'Real Interactive Map',
         'CartoDB Dark Matter tiles — fully zoomable, scrollable, pannable. Built with react-leaflet.'),
        (GRN, '🛡', 'Safe Route (Green Line)',
         'L-shaped path via lit streets, open stores, CCTV coverage. 2.8 km · 34 min.'),
        (RED, '⚠', 'Unsafe Route (Red Dashed)',
         'Direct diagonal shortcut through dim alleys. 2.2 km · 27 min — but flagged as only 38% safe.'),
        (PUR, '📋', 'Route Cards Below Map',
         'Tappable cards with safety %, distance, time, and key hazard warnings for each route.'),
    ]
    for fi, (col, icon, title, body) in enumerate(FEATS10):
        feat_row_clean(sl, RX, 2.55 + fi * 0.84, RW, col, icon, title, body, divider=(fi < len(FEATS10) - 1))

    footer(sl, 10, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 11 — PROTOTYPE: Screen 2 — Navigation
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 4 – Prototype', GRN, 'Screen 2: Active Navigation', 'Turn-by-Turn with Safety Controls')

    LX = 0.55; LW = 9.50
    # WHAT THIS SCREEN DOES
    wtsd = add_rect(sl, LX, 1.50, LW, 0.90, RGBColor(0x00, 0x28, 0x18), GRN, 1.2)
    wtsdtf = wtsd.text_frame; wtsdtf.word_wrap = True
    wtsdtf.margin_left = Inches(0.15); wtsdtf.margin_top = Inches(0.09)
    wp = wtsdtf.paragraphs[0]; wp.text = '🛡  WHAT THIS SCREEN DOES'
    wp.font.name = 'Inter'; wp.font.size = Pt(9.5); wp.font.bold = True; wp.font.color.rgb = GRN; wp.space_after = Pt(4)
    wp2 = wtsdtf.add_paragraph()
    wp2.text = 'The active navigation HUD — turn-by-turn guidance with a safety overlay badge, SOS button always visible, live hazard reporting and one-tap location share.'
    wp2.font.name = 'Inter'; wp2.font.size = Pt(10); wp2.font.color.rgb = T1; wp2.line_spacing = 1.30

    FEATS11 = [
        (GRN,  '🧭', 'Turn-by-Turn HUD',
         'Shows next turn instruction, km remaining, and ETA at top. Updates in real time during walk.'),
        (GRN,  '🛡', 'Safety Corridor Badge',
         '"94% Safety Corridor" overlaid on map — confirms user is on the safe path at all times.'),
        (RED,  '🚨', 'SOS Button (Always Visible)',
         'Animated glowing red button always on screen. One tap starts the 3-second emergency countdown.'),
        (AMB,  '📍', 'Report Hazard Button',
         'Slide-up modal to pin a new danger — dim light, blocked path, or crowd — on the live map.'),
        (BLU,  '📤', 'Share Location Button',
         'One tap sends a live GPS link to pre-set emergency contacts via SMS. Toast confirmation shown.'),
    ]
    for fi, (col, icon, title, body) in enumerate(FEATS11):
        feat_row_clean(sl, LX, 2.55 + fi * 0.76, LW, col, icon, title, body, divider=(fi < len(FEATS11) - 1))

    draw_phone_mockup(sl, 10.50, 1.50, 'navigation')
    cap_tb = add_tb(sl, 10.10, 6.30, 3.0, 0.3)
    para(cap_tb.text_frame, 'Navigation Screen', size=8, color=T2, align=PP_ALIGN.CENTER, add=False)
    footer(sl, 11, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 12 — PROTOTYPE: Screen 3 — SOS
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 4 – Prototype', RED, 'Screen 3: SOS Emergency', '3-Second Countdown & Contact Dispatch')

    draw_phone_mockup(sl, 0.55, 1.50, 'sos')
    cap_tb = add_tb(sl, 0.20, 6.30, 3.0, 0.3)
    para(cap_tb.text_frame, 'SOS Countdown Screen', size=8, color=T2, align=PP_ALIGN.CENTER, add=False)

    RX2 = 3.25; RW2 = 9.50
    # WHAT THIS SCREEN DOES — red banner
    wtsd = add_rect(sl, RX2, 1.50, RW2, 0.90, RGBColor(0x28, 0x08, 0x10), RED, 1.2)
    wtsdtf = wtsd.text_frame; wtsdtf.word_wrap = True
    wtsdtf.margin_left = Inches(0.15); wtsdtf.margin_top = Inches(0.09)
    wp = wtsdtf.paragraphs[0]; wp.text = '🚨  WHAT THIS SCREEN DOES'
    wp.font.name = 'Inter'; wp.font.size = Pt(9.5); wp.font.bold = True; wp.font.color.rgb = RED; wp.space_after = Pt(4)
    wp2 = wtsdtf.add_paragraph()
    wp2.text = 'After the SOS tap, a 3-second animated countdown gives the user time to abort if accidental — then dispatches live GPS via SMS to all emergency contacts instantly.'
    wp2.font.name = 'Inter'; wp2.font.size = Pt(10); wp2.font.color.rgb = T1; wp2.line_spacing = 1.30

    FEATS12 = [
        (RED,  '🔢', '3-Second Visual Countdown',
         'Large pulsing number with expanding ring animation — user sees exactly when GPS is sent.'),
        (AMB,  '✋', 'Hold to Cancel (Abort Button)',
         'Prominent abort button prevents false alarms — key usability insight from early testing.'),
        (GRN,  '✅', 'Live Status Checklist',
         'Shows: GPS lock → SMS compose → Campus security alert in real time as steps complete.'),
        (RED,  '🔊', 'Audio Siren Activated',
         'After 3s, device emits loud alarm siren + dispatches GPS SMS to Mom, Roommate & Campus Police.'),
    ]
    for fi, (col, icon, title, body) in enumerate(FEATS12):
        feat_row_clean(sl, RX2, 2.55 + fi * 0.84, RW2, col, icon, title, body, divider=(fi < len(FEATS12) - 1))

    # Usability insight callout (bottom)
    ut = add_rect(sl, RX2, 6.00, RW2, 0.80, RGBColor(0x28, 0x18, 0x04), AMB, 1.0)
    uttf = ut.text_frame; uttf.word_wrap = True
    uttf.margin_left = Inches(0.15); uttf.margin_top = Inches(0.10)
    up = uttf.paragraphs[0]; up.text = '🧪  USABILITY INSIGHT'
    up.font.name = 'Inter'; up.font.size = Pt(9.5); up.font.bold = True; up.font.color.rgb = AMB; up.space_after = Pt(3)
    up2 = uttf.add_paragraph()
    up2.text = 'Instant SOS (v1) caused anxiety. The 3-second abort window reduced false alarm fear by 100% in re-testing.'
    up2.font.name = 'Inter'; up2.font.size = Pt(9); up2.font.color.rgb = T2; up2.line_spacing = 1.20

    footer(sl, 12, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 13 — PROTOTYPE: Screen 4 — Hazard Report
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 4 – Prototype', AMB, 'Screen 4: Hazard Report', 'Community-Powered Safety — Map Every Danger')

    LX2 = 0.55; LW2 = 9.50
    # WHAT THIS SCREEN DOES — amber banner
    wtsd = add_rect(sl, LX2, 1.50, LW2, 0.90, RGBColor(0x28, 0x18, 0x04), AMB, 1.2)
    wtsdtf = wtsd.text_frame; wtsdtf.word_wrap = True
    wtsdtf.margin_left = Inches(0.15); wtsdtf.margin_top = Inches(0.09)
    wp = wtsdtf.paragraphs[0]; wp.text = '📍  WHAT THIS SCREEN DOES'
    wp.font.name = 'Inter'; wp.font.size = Pt(9.5); wp.font.bold = True; wp.font.color.rgb = AMB; wp.space_after = Pt(4)
    wp2 = wtsdtf.add_paragraph()
    wp2.text = 'A slide-up bottom sheet modal during navigation. Users pick a hazard category — the pin appears instantly on the live map, visible to all SafeRoute users nearby.'
    wp2.font.name = 'Inter'; wp2.font.size = Pt(10); wp2.font.color.rgb = T1; wp2.line_spacing = 1.30

    FEATS13 = [
        (AMB, '🏷', '4 Hazard Categories',
         'Dim Lighting · Blocked Path · Suspicious Crowd · Unsafe Road — selected in one tap.'),
        (GRN, '📌', 'Instant Map Pin',
         'After Submit, an amber ⚠ pin appears on the live map at the reported GPS coordinate immediately.'),
        (BLU, '👥', 'Community-Powered Data',
         'Every report improves the safety algorithm — future routes automatically avoid newly-pinned hazards.'),
        (PUR, '📱', 'Bottom Sheet UX Pattern',
         'Slides up over the map — user never loses context of where they are while reporting the hazard.'),
    ]
    for fi, (col, icon, title, body) in enumerate(FEATS13):
        feat_row_clean(sl, LX2, 2.55 + fi * 0.84, LW2, col, icon, title, body, divider=(fi < len(FEATS13) - 1))

    # Community flywheel callout (bottom)
    fw = add_rect(sl, LX2, 6.00, LW2, 0.80, BG2, PUR, 1.0)
    fwtf = fw.text_frame; fwtf.word_wrap = True
    fwtf.margin_left = Inches(0.15); fwtf.margin_top = Inches(0.10)
    fwp = fwtf.paragraphs[0]; fwp.text = '🔄  The Community Safety Flywheel'
    fwp.font.name = 'Inter'; fwp.font.size = Pt(9.5); fwp.font.bold = True; fwp.font.color.rgb = T0; fwp.space_after = Pt(3)
    fwp2 = fwtf.add_paragraph()
    fwp2.text = 'More users → More reports → Better data → Smarter routes → More users feel safe → More users join. Like Waze for pedestrian night safety.'
    fwp2.font.name = 'Inter'; fwp2.font.size = Pt(9); fwp2.font.color.rgb = T2; fwp2.line_spacing = 1.20

    draw_phone_mockup(sl, 10.50, 1.50, 'hazard')
    cap_tb = add_tb(sl, 10.10, 6.30, 3.0, 0.3)
    para(cap_tb.text_frame, 'Hazard Report Modal', size=8, color=T2, align=PP_ALIGN.CENTER, add=False)
    footer(sl, 13, TOTAL)



    # ════════════════════════════════════════════════
    # SLIDE 14 — TEST: Usability Testing & Iterations
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Phase 5 – Test & Iterate', BLU, 'Usability Testing', 'Test → Feedback → Iterate → Retest')

    # Metrics
    for mi, (v, lb, col) in enumerate([('8', 'Test participants', PUR), ('3', 'Rounds of testing', BLU),
                                        ('100%', 'Task completion\nafter iteration', GRN), ('4.7★', 'Avg satisfaction\nscore', AMB)]):
        stat_box(sl, 0.55 + mi * 3.1, 1.78, 2.85, 1.30, v, lb, col)

    # Iteration table
    ITERS = [
        ('Accidental SOS trigger\nfear during fast walking',
         '7 of 8 testers worried about\nfalse alarm dispatch.',
         'Added 3s countdown with large\n"Hold to Cancel" abort button.',
         'False alarm fear: 100% → 0%', RED),
        ('"38% safety score" not\nunderstood by users',
         '6 of 8 wanted specific reasons,\nnot just a percentage number.',
         'Added ⚠ Dim Alley hazard pin\nlabels on the unsafe route.',
         'Route clarity: 40% → 95%', AMB),
        ('Status bar clipped on\nsmall screens',
         'Phone frame cut off the\n9:41 status bar at top.',
         'CSS min() responsive phone + 100vh\nlocked container in prototype.',
         'Status bar: always visible', BLU),
    ]
    COLS = ['Pain Point', 'Finding', 'Design Iteration', 'Impact']
    COL_COLS = [RED, AMB, BLU, GRN]
    COL_W = [2.7, 2.7, 3.3, 2.7]

    # Header row
    x = 0.55
    for hi, (hdr_text, cw, hcol) in enumerate(zip(COLS, COL_W, COL_COLS)):
        hc = add_rect(sl, x, 3.22, cw, 0.32, BG3, hcol, 0.8)
        htb = add_tb(sl, x+0.08, 3.28, cw-0.16, 0.22)
        para(htb.text_frame, hdr_text, size=9, bold=True, color=hcol, align=PP_ALIGN.CENTER, add=False)
        x += cw + 0.06

    for ri, (pain, finding, iter_, impact, col) in enumerate(ITERS):
        x = 0.55
        for ci, (cell_text, cw) in enumerate(zip([pain, finding, iter_, impact], COL_W)):
            c = add_rect(sl, x, 3.60 + ri * 1.12, cw, 1.05, BG2, COL_COLS[ci] if ci==3 else BDR, 0.6)
            ctf = c.text_frame; ctf.word_wrap = True
            ctf.margin_left = Inches(0.10); ctf.margin_top = Inches(0.08)
            p = ctf.paragraphs[0]; p.text = cell_text
            p.font.name='Inter'; p.font.size=Pt(9.5)
            p.font.color.rgb = GRN if ci==3 else T1; p.font.bold = ci==3; p.line_spacing=1.2
            x += cw + 0.06

    footer(sl, 14, TOTAL)

    # ════════════════════════════════════════════════
    # SLIDE 15 — CONCLUSION
    # ════════════════════════════════════════════════
    sl = prs.slides.add_slide(blank)
    solid_bg(sl)
    slide_header(sl, 'Conclusion', GRN, 'SafeRoute — Delivered', 'From Research to Working Prototype in 5 HCD Phases')

    # Left: Phase summary (2x3 grid)
    PHASES = [
        ('01 Discover', BLU, 'Speed-bias in nav apps. 73% women anxious at night. 85% prefer longer lit routes.', 0.55, 1.50, BG2),
        ('02 Define', PUR, 'Elena persona, Empathy Map, User Journey. Fear → Relief arc clearly mapped.', 5.50, 1.50, BG2),
        ('03 Ideate', AMB, '6 Crazy 8 features + 3 HMW statements directly mapped to built features.', 0.55, 2.50, BG2),
        ('04 Design', GRN, 'Dark design system, 4-screen IA, design tokens, component library.', 5.50, 2.50, BG2),
        ('05 Test', RED, '3 rounds, 8 users. SOS abort + hazard pins from testing. 100% completion.', 0.55, 3.50, BG2),
        ('✦ Live Prototype', GRN, 'Real Leaflet map · Haversine distances · All 4 screens · SOS flow · Hazard reports · Fully functional', 5.50, 3.50, RGBColor(0x00, 0x1A, 0x0A)),
    ]
    for phase, col, summary, sx, sy, bg in PHASES:
        pc = add_rect(sl, sx, sy, 4.65, 0.88, bg, col, 0.8 if bg == BG2 else 1.0)
        pctf = pc.text_frame; pctf.word_wrap=True
        pctf.margin_left=Inches(0.12); pctf.margin_top=Inches(0.06)
        pp1 = pctf.paragraphs[0]; pp1.text = phase
        pp1.font.name='Inter'; pp1.font.size=Pt(10.5); pp1.font.bold=True; pp1.font.color.rgb=col; pp1.space_after=Pt(2)
        pp2 = pctf.add_paragraph(); pp2.text = summary
        pp2.font.name='Inter'; pp2.font.size=Pt(9); pp2.font.color.rgb=T2; pp2.line_spacing=1.1

    # Bottom checklist card
    wwd_card = add_rect(sl, 0.55, 4.50, 9.60, 2.40, BG2, BDR, 0.6)
    wwd_tb = add_tb(sl, 0.70, 4.58, 9.30, 0.30)
    para(wwd_tb.text_frame, 'WHAT WAS DELIVERED', size=10, bold=True, color=T2, add=False)

    col1_tb = add_tb(sl, 0.70, 4.92, 4.50, 1.85)
    c1_tf = col1_tb.text_frame; c1_tf.word_wrap=True
    DEL1 = [
        '🗺  Real interactive Leaflet map (react-leaflet + CartoDB)',
        '📏  Actual route distances via Haversine formula',
        '🛡  Safe vs unsafe route comparison with safety %',
        '🚨  Full SOS countdown → contact dispatch flow',
    ]
    for item in DEL1:
        p = c1_tf.add_paragraph() if c1_tf.paragraphs[0].text else c1_tf.paragraphs[0]
        p.text = f'✓  {item}'; p.space_after = Pt(4); p.line_spacing = 1.15
        p.font.name='Inter'; p.font.size=Pt(9.5); p.font.color.rgb=T1

    col2_tb = add_tb(sl, 5.50, 4.92, 4.50, 1.85)
    c2_tf = col2_tb.text_frame; c2_tf.word_wrap=True
    DEL2 = [
        '📍  Community hazard report → live map pin',
        '📤  Location sharing with toast notifications',
        '🌙  4-screen native app with dark UI system',
        '📱  Full phone frame with status bar + home bar',
    ]
    for item in DEL2:
        p = c2_tf.add_paragraph() if c2_tf.paragraphs[0].text else c2_tf.paragraphs[0]
        p.text = f'✓  {item}'; p.space_after = Pt(4); p.line_spacing = 1.15
        p.font.name='Inter'; p.font.size=Pt(9.5); p.font.color.rgb=T1

    # Right: Phone mockup
    draw_phone_mockup(sl, 10.45, 1.50, 'dashboard')
    
    # Bottom Right button card
    btn = add_rect(sl, 10.45, 6.15, 2.30, 0.75, RGBColor(0x00, 0x1A, 0x0A), GRN, 1.0)
    btn_tb = add_tb(sl, 10.45, 6.22, 2.30, 0.60)
    tf_btn = btn_tb.text_frame; tf_btn.word_wrap = True
    para(tf_btn, 'Tab 3', size=10, bold=True, color=GRN, align=PP_ALIGN.CENTER, add=False)
    para(tf_btn, 'Live Prototype →', size=9, color=T1, align=PP_ALIGN.CENTER, add=True)

    footer(sl, 15, TOTAL)

    # Try to save to v5 version to avoid locked file errors
    out = 'SafeRoute_HCD_Presentation_v5.pptx'
    try:
        prs.save(out)
        print(f'Saved: {out}')
    except PermissionError:
        print(f'Warning: Could not save {out} due to PermissionError (file locked)')

    try:
        prs.save('SafeRoute_HCD_Presentation.pptx')
        print('Saved: SafeRoute_HCD_Presentation.pptx')
    except PermissionError:
        print('Warning: Could not save SafeRoute_HCD_Presentation.pptx due to PermissionError (file locked)')

    # Save to public folder so browser downloads it
    os.makedirs('public', exist_ok=True)
    try:
        prs.save('public/SafeRoute_HCD_Presentation.pptx')
        print('Saved: public/SafeRoute_HCD_Presentation.pptx')
    except PermissionError:
        print('Warning: Could not save public/SafeRoute_HCD_Presentation.pptx due to PermissionError (file locked)')

    return out


if __name__ == '__main__':
    build()
